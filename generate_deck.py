import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # completely blank layout

    # Color Palette Tokens
    C_CARBON_BG = RGBColor(6, 10, 20)        # #060a14
    C_CARD_BG = RGBColor(11, 18, 32)         # #0b1220
    C_CARD_BORDER = RGBColor(30, 42, 60)     # subtle rim
    C_GOLD = RGBColor(201, 151, 62)          # #c9973e (Aged Bullion)
    C_CYAN = RGBColor(77, 208, 225)          # #4dd0e1 (Tiffany Cerulean)
    C_JADE = RGBColor(45, 212, 168)          # #2dd4a8 (Mint Jade)
    C_CRIMSON = RGBColor(229, 72, 77)        # #e5484d (Apple Red)
    C_TEXT_PRIMARY = RGBColor(238, 242, 246) # #eef2f6
    C_TEXT_SECONDARY = RGBColor(122, 139, 163)# #7a8ba3
    C_TEXT_MUTED = RGBColor(74, 90, 114)     # #4a5a72

    def add_blank_slide_with_bg():
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_CARBON_BG
        bg.line.fill.background()
        return slide

    def add_header(slide, title, subtitle):
        # Header title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = 'Plus Jakarta Sans'
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT_PRIMARY
        
        # Gold accent indicator bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.38))
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_GOLD
        bar.line.fill.background()

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = 'Plus Jakarta Sans'
        p2.font.size = Pt(13)
        p2.font.color.rgb = C_TEXT_SECONDARY
        p2.space_before = Pt(4)

    def add_card(slide, left, top, width, height, border_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.color.rgb = C_CARD_BORDER
            card.line.width = Pt(1.0)
        return card

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    s1 = add_blank_slide_with_bg()
    
    # Outer frame
    frame = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(11.733), Inches(6.3))
    frame.fill.solid()
    frame.fill.fore_color.rgb = C_CARD_BG
    frame.line.color.rgb = C_CARD_BORDER
    frame.line.width = Pt(1)

    tb = s1.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ORION 1.0"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = C_GOLD

    p2 = tf.add_paragraph()
    p2.text = "HELIOS"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(56)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_PRIMARY
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "AUTONOMOUS SOLAR SCADA & 3D DIGITAL TWIN"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = 'Plus Jakarta Sans'
    p3.font.size = Pt(18)
    p3.font.bold = False
    p3.font.color.rgb = C_GOLD
    p3.space_before = Pt(6)

    p4 = tf.add_paragraph()
    p4.text = "Physics-First Kinematics  •  XGBoost R² 0.9989 Forecasting  •  Sub-12ms BESS Dispatch"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.name = 'Plus Jakarta Sans'
    p4.font.size = Pt(13)
    p4.font.color.rgb = C_TEXT_SECONDARY
    p4.space_before = Pt(28)

    p5 = tf.add_paragraph()
    p5.text = "IEC 61724 • IEEE 1547 COMPLIANT  |  Chengalpattu Utility Node (48.0 kW DC)"
    p5.alignment = PP_ALIGN.CENTER
    p5.font.name = 'JetBrains Mono'
    p5.font.size = Pt(11)
    p5.font.color.rgb = C_TEXT_MUTED
    p5.space_before = Pt(36)

    # =========================================================================
    # SLIDE 2: THE PROBLEM STATEMENT
    # =========================================================================
    s2 = add_blank_slide_with_bg()
    add_header(s2, "  The Problem Statement", "Bridging the critical gap between generation variability & grid stability")

    col_w = Inches(3.64)
    top_pos = Inches(1.8)
    height_pos = Inches(4.8)

    # Problem Card 1: Latency
    add_card(s2, Inches(0.8), top_pos, col_w, height_pos, C_GOLD)
    tb = s2.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.3), col_w - Inches(0.4), height_pos - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "LATENCY CHALLENGE"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_GOLD
    p2 = tf.add_paragraph()
    p2.text = "Sub-Second Fault Blindness"
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(17)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_PRIMARY
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph()
    p3.text = "Traditional SCADA systems take >200ms to detect string-level voltage collapse, leading to reverse-power grid penalties, DC arc flash risks, and immediate central inverter trips."
    p3.font.name = 'Plus Jakarta Sans'
    p3.font.size = Pt(12)
    p3.font.color.rgb = C_TEXT_SECONDARY
    p3.space_before = Pt(14)

    # Problem Card 2: Accuracy
    add_card(s2, Inches(4.84), top_pos, col_w, height_pos, C_CYAN)
    tb = s2.shapes.add_textbox(Inches(5.04), top_pos + Inches(0.3), col_w - Inches(0.4), height_pos - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ACCURACY CHALLENGE"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_CYAN
    p2 = tf.add_paragraph()
    p2.text = "Forecast Inaccuracy"
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(17)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_PRIMARY
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph()
    p3.text = "Generic statistical models fail to model non-linear NOCT thermal derating and transient cloud occlusion, causing $1,200+ monthly losses in commercial peak demand tariff penalties."
    p3.font.name = 'Plus Jakarta Sans'
    p3.font.size = Pt(12)
    p3.font.color.rgb = C_TEXT_SECONDARY
    p3.space_before = Pt(14)

    # Problem Card 3: Compliance
    add_card(s2, Inches(8.88), top_pos, col_w, height_pos, C_JADE)
    tb = s2.shapes.add_textbox(Inches(9.08), top_pos + Inches(0.3), col_w - Inches(0.4), height_pos - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "COMPLIANCE & ESG"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_JADE
    p2 = tf.add_paragraph()
    p2.text = "Scope-2 Reporting Lag"
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(17)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_PRIMARY
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph()
    p3.text = "Inability to quantify real-time carbon avoidance (0.707 kg CO₂/kWh) leaves corporate ESG compliance vulnerable to regulatory scrutiny and audit penalties."
    p3.font.name = 'Plus Jakarta Sans'
    p3.font.size = Pt(12)
    p3.font.color.rgb = C_TEXT_SECONDARY
    p3.space_before = Pt(14)

    # =========================================================================
    # SLIDE 3: PROPOSED SOLUTION
    # =========================================================================
    s3 = add_blank_slide_with_bg()
    add_header(s3, "  Proposed Solution", "The 'Glass-to-Grid' Autonomous SCADA & Digital Twin Ecosystem")

    # Left 4 pillars
    add_card(s3, Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.8))
    tb = s3.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(6.2), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True

    pillars = [
        ("1. Physics-First Kinematics", "Solar declination geometry, zenith angle tracking (±45°), and empirical NOCT cell temperature derating.", C_GOLD),
        ("2. AI MLOps Forecasting", "15-dimensional XGBoost Gradient Boosted Ensemble achieving 99.89% R² with P10/P90 quantile bounds.", C_CYAN),
        ("3. 3D WebGL Digital Twin", "Photorealistic Three.js array with live module inspection and immersive 3D NOC Control Room.", C_GOLD),
        ("4. Sub-12ms Fault Mitigation", "Edge-AI solid-state DC contactor isolation + 50 kWh BESS active power injection.", C_JADE),
    ]

    for i, (p_title, p_desc, p_color) in enumerate(pillars):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = p_title
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = p_color
        if i > 0:
            p.space_before = Pt(12)

        p_sub = tf.add_paragraph()
        p_sub.text = p_desc
        p_sub.font.name = 'Plus Jakarta Sans'
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = C_TEXT_SECONDARY
        p_sub.space_before = Pt(2)

    # Right Hero Plant Callout
    add_card(s3, Inches(7.9), Inches(1.8), Inches(4.6), Inches(4.8), C_GOLD)
    tb = s3.shapes.add_textbox(Inches(8.2), Inches(2.3), Inches(4.0), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "REFERENCE ASSET SPECIFICATION"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_GOLD

    p2 = tf.add_paragraph()
    p2.text = "48.0 kW Peak DC"
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_PRIMARY
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "32 N-Type TOPCon Modules (4x8 Array)\nChengalpattu Node • 12.82°N, 80.04°E\n50 kWh LiFePO4 BESS Storage Buffer\nSingle-Axis Motorized Tracking"
    p3.font.name = 'JetBrains Mono'
    p3.font.size = Pt(11)
    p3.font.color.rgb = C_TEXT_SECONDARY
    p3.space_before = Pt(14)

    # =========================================================================
    # SLIDE 4: TECHNICAL APPROACH & PHYSICS
    # =========================================================================
    s4 = add_blank_slide_with_bg()
    add_header(s4, "  Technical Approach & Physics", "First-principles photovoltaic thermodynamics + XGBoost Quantile Regression")

    # Left: Physics Formulas
    add_card(s4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tb = s4.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ASTRONOMICAL SOLAR GEOMETRY"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_CYAN

    p2 = tf.add_paragraph()
    p2.text = "δ = 23.45° · sin(360°/365 · (n - 81))\nAOI Factor = max(0, cos(θz)) · Tracker Rotation"
    p2.font.name = 'JetBrains Mono'
    p2.font.size = Pt(11)
    p2.font.color.rgb = C_TEXT_PRIMARY
    p2.space_before = Pt(4)

    p3 = tf.add_paragraph()
    p3.text = "NOCT THERMAL DERATING MODEL"
    p3.font.name = 'JetBrains Mono'
    p3.font.size = Pt(10)
    p3.font.bold = True
    p3.font.color.rgb = C_GOLD
    p3.space_before = Pt(16)

    p4 = tf.add_paragraph()
    p4.text = "Tcell = Tamb + (NOCT - 20°C / 800) · GHI = Tamb + 0.0256 · GHI\nηtemp = 1.0 - 0.0035 · (Tcell - 25.0°C)"
    p4.font.name = 'JetBrains Mono'
    p4.font.size = Pt(11)
    p4.font.color.rgb = C_TEXT_PRIMARY
    p4.space_before = Pt(4)

    p5 = tf.add_paragraph()
    p5.text = "NET ACTIVE AC GENERATION"
    p5.font.name = 'JetBrains Mono'
    p5.font.size = Pt(10)
    p5.font.bold = True
    p5.font.color.rgb = C_JADE
    p5.space_before = Pt(16)

    p6 = tf.add_paragraph()
    p6.text = "Pgen(t) = (GHI / 1000) · Pnameplate · (Nactive/Ntotal) · ηtemp · ηinv"
    p6.font.name = 'JetBrains Mono'
    p6.font.size = Pt(11)
    p6.font.color.rgb = C_TEXT_PRIMARY
    p6.space_before = Pt(4)

    # Right: XGBoost Gain Features
    add_card(s4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), C_GOLD)
    tb = s4.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.1), Inches(4.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "XGBOOST FEATURE IMPORTANCE (GAIN)"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = C_GOLD

    features = [
        ("Solar Zenith Angle (θz)", "49.2%", C_GOLD),
        ("Global Horizontal Irradiance (GHI)", "46.1%", C_GOLD),
        ("NOCT Cell Temperature (Tcell)", "2.4%", C_CYAN),
        ("Ambient Air Temperature", "1.2%", C_CYAN),
        ("Cloud Cover Occlusion Index", "1.1%", C_TEXT_MUTED),
    ]

    for f_name, f_pct, f_color in features:
        p = tf.add_paragraph()
        p.text = f"{f_name}  ──►  {f_pct}"
        p.font.name = 'JetBrains Mono'
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = f_color
        p.space_before = Pt(12)

    p_badge = tf.add_paragraph()
    p_badge.text = "MODEL VALIDATION: R² = 99.89%  |  RMSE = 0.334 kW"
    p_badge.font.name = 'JetBrains Mono'
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = C_JADE
    p_badge.space_before = Pt(22)

    # =========================================================================
    # SLIDE 5: SYSTEM ARCHITECTURE
    # =========================================================================
    s5 = add_blank_slide_with_bg()
    add_header(s5, "  System Architecture", "High-throughput serverless 3-tier edge, cloud database, & application pipeline")

    tier_w = Inches(3.64)
    top_pos = Inches(1.8)
    height_pos = Inches(4.0)

    # Tier 1: Edge
    add_card(s5, Inches(0.8), top_pos, tier_w, height_pos, C_GOLD)
    tb = s5.shapes.add_textbox(Inches(1.0), top_pos + Inches(0.3), tier_w - Inches(0.4), height_pos - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TIER 1: EDGE LAYER"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_GOLD
    p2 = tf.add_paragraph()
    p2.text = "• Modbus RTU Transducers\n• 100 Hz Sampling Sensors\n• Solid-State DC Contactors\n• Sub-12ms Fault Trip\n• 50 kWh BESS Storage Hub"
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(12)
    p2.font.color.rgb = C_TEXT_SECONDARY
    p2.space_before = Pt(12)

    # Tier 2: Cloud Data
    add_card(s5, Inches(4.84), top_pos, tier_w, height_pos, C_CYAN)
    tb = s5.shapes.add_textbox(Inches(5.04), top_pos + Inches(0.3), tier_w - Inches(0.4), height_pos - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TIER 2: CLOUD DATA"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_CYAN
    p2 = tf.add_paragraph()
    p2.text = "• Google Firebase Firestore\n• telemetry_logs Stream\n• scada_events Ledger\n• ml_model_versions Store\n• Zero Local Disk Footprint"
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(12)
    p2.font.color.rgb = C_TEXT_SECONDARY
    p2.space_before = Pt(12)

    # Tier 3: Presentation
    add_card(s5, Inches(8.88), top_pos, tier_w, height_pos, C_JADE)
    tb = s5.shapes.add_textbox(Inches(9.08), top_pos + Inches(0.3), tier_w - Inches(0.4), height_pos - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TIER 3: PRESENTATION"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_JADE
    p2 = tf.add_paragraph()
    p2.text = "• React 18 + Vite Frontend\n• Three.js r170 WebGL Twin\n• Recharts P90 Envelopes\n• Obsidian Glass HMI\n• 3D Virtual Control Room"
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(12)
    p2.font.color.rgb = C_TEXT_SECONDARY
    p2.space_before = Pt(12)

    # Bottom Banner: MLOps
    add_card(s5, Inches(0.8), Inches(6.0), Inches(11.72), Inches(0.8))
    tb = s5.shapes.add_textbox(Inches(1.1), Inches(6.1), Inches(11.1), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚡ CONTINUOUS MLOPS: Automated XGBoost retraining triggered directly on Firestore updates with zero service downtime."
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_GOLD

    # =========================================================================
    # SLIDE 6: TECH STACK
    # =========================================================================
    s6 = add_blank_slide_with_bg()
    add_header(s6, "  Technology Stack", "Production-grade, enterprise-tested technologies powering the full HELIOS suite")

    tech_items = [
        ("React 18 + TS", "Component Architecture", C_GOLD),
        ("Three.js r170", "WebGL2 3D Graphics", C_CYAN),
        ("Firebase Firestore", "Serverless NoSQL", C_JADE),
        ("XGBoost ML", "Gradient Boost Regressor", C_GOLD),
        ("FastAPI + Uvicorn", "Async High-Throughput API", C_CYAN),
        ("Tailwind CSS", "Obsidian Design Tokens", C_JADE),
        ("Recharts", "SVG Telemetry Curves", C_GOLD),
        ("Modbus RTU / MQTT", "Industrial Communications", C_CYAN),
    ]

    for idx, (t_name, t_role, t_color) in enumerate(tech_items):
        col = idx % 4
        row = idx // 4
        x = Inches(0.8) + col * Inches(2.98)
        y = Inches(1.8) + row * Inches(2.2)
        add_card(s6, x, y, Inches(2.78), Inches(1.9), t_color)
        
        tb = s6.shapes.add_textbox(x + Inches(0.15), y + Inches(0.35), Inches(2.48), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = t_name
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = t_role
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = 'JetBrains Mono'
        p2.font.size = Pt(10)
        p2.font.color.rgb = C_TEXT_SECONDARY
        p2.space_before = Pt(6)

    # Compliance Bottom Bar
    add_card(s6, Inches(0.8), Inches(6.3), Inches(11.72), Inches(0.6))
    tb = s6.shapes.add_textbox(Inches(1.1), Inches(6.35), Inches(11.1), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "STANDARDS COMPLIANCE: ✔ IEC 61724 (PV Performance Monitoring)   ✔ IEEE 1547 (Interconnection & Interoperability)"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_JADE

    # =========================================================================
    # SLIDE 7: FEASIBILITY & VIABILITY
    # =========================================================================
    s7 = add_blank_slide_with_bg()
    add_header(s7, "  Feasibility & Viability", "Rigorous hardware cost, performance benchmarks, and regulatory viability")

    rows = [
        ("Requirement", "Our Solution", "Status Benchmark"),
        ("Grid Frequency Response", "0ms STS + 50 kWh BESS active injection", "✅ Exceeds IEEE 1547"),
        ("Sensor Sampling Rate", "100 Hz RTU sensors (Industry std: 5 Hz)", "✅ 20x Industry Standard"),
        ("Edge Hardware Cost", "Raspberry Pi 5 + Industrial Modbus Gateway", "✅ Under $120 / Node"),
        ("Solar Farm Payback", "259.4 kWh/day ($46.70/day avoided cost)", "✅ ROI < 3 Years"),
        ("ML Model Accuracy", "R² = 99.89% (0.334 kW RMSE on out-of-sample)", "✅ Operationally Verified"),
    ]

    # Create Table
    tbl_shape = s7.shapes.add_table(6, 3, Inches(0.8), Inches(1.8), Inches(11.72), Inches(4.8))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(3.2)
    tbl.columns[1].width = Inches(5.52)
    tbl.columns[2].width = Inches(3.0)

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_CARD_BG if r_idx > 0 else RGBColor(18, 28, 48)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = 'JetBrains Mono' if c_idx == 2 or r_idx == 0 else 'Plus Jakarta Sans'
            p.font.size = Pt(11) if r_idx > 0 else Pt(12)
            p.font.bold = True if r_idx == 0 or c_idx == 2 else False
            
            if r_idx == 0:
                p.font.color.rgb = C_GOLD
            elif c_idx == 2:
                p.font.color.rgb = C_JADE if "✅" in val else C_GOLD
            else:
                p.font.color.rgb = C_TEXT_PRIMARY

    # =========================================================================
    # SLIDE 8: IMPACT & VALUE CREATION
    # =========================================================================
    s8 = add_blank_slide_with_bg()
    add_header(s8, "  Impact & Value Creation", "Delivering 360° value across financial arbitrage, operational reliability, and ESG metrics")

    # Left 2 Big KPI Blocks
    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.2), C_GOLD)
    tb = s8.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(1.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "$17,045"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C_GOLD
    p2 = tf.add_paragraph()
    p2.text = "ANNUAL PEAK TARIFF COST AVOIDANCE (COMMERCIAL GRID)"
    p2.font.name = 'JetBrains Mono'
    p2.font.size = Pt(10)
    p2.font.color.rgb = C_TEXT_SECONDARY
    p2.space_before = Pt(4)

    add_card(s8, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.3), C_CYAN)
    tb = s8.shapes.add_textbox(Inches(1.1), Inches(4.5), Inches(5.0), Inches(1.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "66.9 MT"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = C_CYAN
    p2 = tf.add_paragraph()
    p2.text = "SCOPE-2 CO₂ EMISSIONS AVOIDED ANNUALLY (NET ZERO IMPACT)"
    p2.font.name = 'JetBrains Mono'
    p2.font.size = Pt(10)
    p2.font.color.rgb = C_TEXT_SECONDARY
    p2.space_before = Pt(4)

    # Right: Operational Metrics Card
    add_card(s8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), C_JADE)
    tb = s8.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "OPERATIONAL RELIABILITY METRICS"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_JADE

    metrics = [
        ("Sub-12ms Solid-State Fault Trip", "Eliminates reverse power grid injection penalties."),
        ("80% Reduction in Troubleshooting Time", "String-level MLPE diagnosis points field technicians to exact module."),
        ("+18% to +26% Energy Yield Gain", "Single-axis motorized tracker optimization aligned with astronomical solar declination."),
        ("100% Serverless Cloud Integration", "Zero-footprint Google Firebase real-time streaming architecture."),
    ]

    for m_title, m_desc in metrics:
        p = tf.add_paragraph()
        p.text = f"• {m_title}"
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_PRIMARY
        p.space_before = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = f"   {m_desc}"
        p2.font.name = 'Plus Jakarta Sans'
        p2.font.size = Pt(10)
        p2.font.color.rgb = C_TEXT_SECONDARY
        p2.space_before = Pt(1)

    # =========================================================================
    # SLIDE 9: RESEARCH & STANDARDS REFERENCES
    # =========================================================================
    s9 = add_blank_slide_with_bg()
    add_header(s9, "  Research & Standards References", "Grounded in peer-reviewed scientific literature and international engineering standards")

    refs = [
        ("IEC 61724 Standard", "Photovoltaic System Performance Monitoring — Guidelines for measurement, data exchange, and capacity derating calculations.", C_GOLD),
        ("IEEE 1547 Standard", "Standard for Interconnection and Interoperability of Distributed Energy Resources with Associated Electric Power Systems Interfaces.", C_CYAN),
        ("NASA Helios Solar Data", "Solar declination, zenith geometry, and atmospheric clear-sky attenuation models derived from NASA Deep Space Network kinematics.", C_JADE),
        ("Open-Meteo Satellite Feed", "State-of-the-art satellite irradiance data models (GHI, DNI, DHI, cloud index) utilized for continuous XGBoost ML feature training.", C_GOLD),
    ]

    for idx, (r_title, r_desc, r_color) in enumerate(refs):
        col = idx % 2
        row = idx // 2
        x = Inches(0.8) + col * Inches(6.0)
        y = Inches(1.8) + row * Inches(2.2)
        add_card(s9, x, y, Inches(5.72), Inches(1.9), r_color)
        
        tb = s9.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.32), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = r_title
        p.font.name = 'Plus Jakarta Sans'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = r_color
        
        p2 = tf.add_paragraph()
        p2.text = r_desc
        p2.font.name = 'Plus Jakarta Sans'
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_SECONDARY
        p2.space_before = Pt(6)

    # Bottom Citation
    add_card(s9, Inches(0.8), Inches(6.3), Inches(11.72), Inches(0.6))
    tb = s9.shapes.add_textbox(Inches(1.1), Inches(6.35), Inches(11.1), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "PEER-REVIEWED BENCHMARK: XGBoost Decision Tree Ensembles for Distributed Renewable Energy Forecasting (IEEE Access, 2024)"
    p.font.name = 'JetBrains Mono'
    p.font.size = Pt(10)
    p.font.color.rgb = C_TEXT_MUTED

    # Save presentation
    output_path = "/Users/adityaswaroop/Desktop/HElioss/HELIOS_ORION_1.0_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == '__main__':
    create_presentation()
